from langchain_core.tools import BaseTool
from typing import Optional
from langchain.callbacks.manager import CallbackManagerForToolRun
import fitz
import docx
import json
from openpyxl import load_workbook
from paddleocr import PaddleOCR
from pptx import Presentation
import os


class DocTool(BaseTool):
    name: str = "doc_tool"
    description: str = """当需要识别文本文件，且文件格式属于[txt，html，xml、docx、xlsx、pptx、csv、json]这些常见的文档类型时，可以通过本工具从文件中提取文本，
    该工具输入文件名称，输出对应文件文本内容"""

    def _run(
        self,
        query: str,
        run_manager: Optional[CallbackManagerForToolRun] = None,
    ) -> str:
        try:
            path = os.path.join(os.getcwd(), "uploaded_files")
            file = os.path.join(path, query)
            text_data = ""
            if query.lower().endswith(".pdf"):
                pdf_data = fitz.open(file)
                for page in pdf_data.pages():
                    if not page.get_text():
                        ocr = PaddleOCR(use_angle_cls=True, lang="ch")
                        result = ocr.ocr(file, cls=True)
                        for idx in range(len(result)):
                            res = result[idx]
                            text_data += res[0][1][0]
                        return text_data
                    else:
                        temp = page.get_text()
                        text_data += temp
            elif query.lower().endswith(".docx"):
                doc = docx.Document(file)
                TextList = []
                for paragraph in doc.paragraphs:
                    TextList.append(paragraph.text)
                text_data = "\n".join(TextList)
            elif query.lower().endswith(".xlsx"):
                with open(file, "rb") as f:
                    workbook = load_workbook(f)
                    if workbook is not None:
                        sheet = workbook.active
                        for row in sheet.iter_rows(values_only=True):
                            text_data += str(row)
            elif query.lower().endswith(".json"):
                with open(file, "r") as f:
                    text_data = str(json.load(f))
            elif query.lower().endswith(".pptx"):
                presentation = Presentation(file)
                for slide in presentation.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "text"):
                            if shape.text:
                                text_data += shape.text
            else:
                f = open(file, encoding="utf-8")
                text_data = f.readlines()
            return text_data
        except FileNotFoundError:
            return "File not found"
        except Exception as e:
            return f"An error occurred: {e}"
